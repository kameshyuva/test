from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Use a blank slide layout with a title (Layout 5)
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)

# Set the title
title = slide.shapes.title
title.text = "Local CPU-Optimized Agent Guardrails Architecture"
title.text_frame.paragraphs[0].font.size = Pt(32)

# ---------------------------------------------------------
# Left Column: Compute, Context, and Deployment
# ---------------------------------------------------------
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
tf_left = left_box.text_frame
tf_left.word_wrap = True

# 1. Compute & Resource
p = tf_left.add_paragraph()
p.text = "1. Compute & Resource Limits"
p.font.bold = True
p.font.size = Pt(18)

p = tf_left.add_paragraph(); p.text = "• Concurrency Locks: asyncio.Semaphore"; p.level = 1; p.font.size = Pt(14)
p = tf_left.add_paragraph(); p.text = "• API Rate Limiting: slowapi"; p.level = 1; p.font.size = Pt(14)
p = tf_left.add_paragraph(); p.text = "• Thread/Token Limits: Ollama (num_thread, num_predict)"; p.level = 1; p.font.size = Pt(14)

# 2. Context & State
p = tf_left.add_paragraph()
p.text = "\n2. Context & State Management"
p.font.bold = True
p.font.size = Pt(18)

p = tf_left.add_paragraph(); p.text = "• Pre-Computation Counting: tiktoken"; p.level = 1; p.font.size = Pt(14)
p = tf_left.add_paragraph(); p.text = "• Asymmetric Compression: qwen3.5:0.5b"; p.level = 1; p.font.size = Pt(14)
p = tf_left.add_paragraph(); p.text = "• Persistent Facts: LlamaIndex Memory + ScyllaDB"; p.level = 1; p.font.size = Pt(14)

# 3. Deployment
p = tf_left.add_paragraph()
p.text = "\n3. Deployment Strategy"
p.font.bold = True
p.font.size = Pt(18)

p = tf_left.add_paragraph(); p.text = "• Image & Dependencies: Docker, uv"; p.level = 1; p.font.size = Pt(14)

# ---------------------------------------------------------
# Right Column: Input and Output Sanitization
# ---------------------------------------------------------
right_box = slide.shapes.add_textbox(Inches(5.0), Inches(1.5), Inches(4.5), Inches(5.5))
tf_right = right_box.text_frame
tf_right.word_wrap = True

# 4. Input Sanitization
p = tf_right.add_paragraph()
p.text = "4. Input Sanitization"
p.font.bold = True
p.font.size = Pt(18)

p = tf_right.add_paragraph(); p.text = "• DoS/Buffer Limit: Python len()"; p.level = 1; p.font.size = Pt(14)
p = tf_right.add_paragraph(); p.text = "• Injection/Obfuscation: Python re (Regex)"; p.level = 1; p.font.size = Pt(14)
p = tf_right.add_paragraph(); p.text = "• PII Redaction: Presidio + spaCy (en_core_web_sm)"; p.level = 1; p.font.size = Pt(14)
p = tf_right.add_paragraph(); p.text = "• Profanity Filter: better_profanity"; p.level = 1; p.font.size = Pt(14)

# 5. Output Sanitization
p = tf_right.add_paragraph()
p.text = "\n5. Output Sanitization"
p.font.bold = True
p.font.size = Pt(18)

p = tf_right.add_paragraph(); p.text = "• XSS/HTML Stripping: bleach"; p.level = 1; p.font.size = Pt(14)
p = tf_right.add_paragraph(); p.text = "• Brand/PII Suppression: Presidio (Custom Recognizer)"; p.level = 1; p.font.size = Pt(14)
p = tf_right.add_paragraph(); p.text = "• Real-Time Leak Checks: Regex Tripwire + qwen3.5:0.5b"; p.level = 1; p.font.size = Pt(14)
p = tf_right.add_paragraph(); p.text = "• Memory Optimization: Ollama keep_alive: -1"; p.level = 1; p.font.size = Pt(14)

# Save the slide
output_path = 'One_Page_Guardrails_Summary.pptx'
prs.save(output_path)
print(f"Slide successfully saved to {output_path}")
